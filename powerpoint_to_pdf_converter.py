
#!/usr/bin/env python3
"""
Enhanced PowerPoint to PDF Converter
Converts PowerPoint presentations to PDF while preserving formatting, layout, images, and slide templates.
Uses multiple methods for optimal conversion quality with advanced formatting preservation.
"""

import os
import tempfile
import traceback
import shutil
from pathlib import Path
from typing import Dict, Any, Optional, List
import uuid
import subprocess
import platform

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    print("python-pptx successfully imported")
except ImportError as e:
    print(f"Warning: python-pptx not available - {e}")
    Presentation = None

try:
    import pypandoc
    print("pypandoc successfully imported")
except ImportError as e:
    print(f"Warning: pypandoc not available - {e}")
    pypandoc = None

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas
    print("reportlab successfully imported")
except ImportError as e:
    print(f"Warning: reportlab not available - {e}")
    SimpleDocTemplate = None

try:
    from PIL import Image
    print("Pillow successfully imported")
except ImportError as e:
    print(f"Warning: Pillow not available - {e}")
    Image = None

class PowerPointToPDFConverter:
    """Enhanced PowerPoint to PDF converter with advanced formatting and content preservation"""
    
    def __init__(self, temp_dir: str = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()
        
    def convert_powerpoint_to_pdf(self, pptx_path: str, output_path: str = None) -> Dict[str, Any]:
        """
        Convert PowerPoint to PDF with enhanced formatting preservation
        
        Args:
            pptx_path: Path to input PowerPoint file
            output_path: Path for output PDF file (optional)
            
        Returns:
            Dictionary with conversion results
        """
        try:
            if not os.path.exists(pptx_path):
                return {"success": False, "error": f"PowerPoint file not found: {pptx_path}"}
            
            # Generate output path if not provided
            if not output_path:
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.pdf")
            
            print(f"Starting PowerPoint to PDF conversion: {pptx_path} -> {output_path}")
            
            # Method 1: LibreOffice conversion (best quality)
            if self._convert_with_libreoffice(pptx_path, output_path):
                print("LibreOffice conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "filename": os.path.basename(output_path),
                    "method": "libreoffice_headless",
                    "message": "Conversion completed with LibreOffice (high quality)",
                    "metadata": self._get_pdf_metadata(output_path)
                }
            
            # Method 2: Direct PPTX parsing with ReportLab
            if self._convert_with_reportlab_method(pptx_path, output_path):
                print("ReportLab conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "filename": os.path.basename(output_path),
                    "method": "reportlab_parsing",
                    "message": "Conversion completed with ReportLab parsing",
                    "metadata": self._get_pdf_metadata(output_path)
                }
            
            # Method 3: Pypandoc conversion (fallback)
            if self._convert_with_pypandoc(pptx_path, output_path):
                print("Pypandoc conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "filename": os.path.basename(output_path),
                    "method": "pypandoc",
                    "message": "Conversion completed with Pypandoc",
                    "metadata": self._get_pdf_metadata(output_path)
                }
            
            return {"success": False, "error": "All conversion methods failed"}
            
        except Exception as e:
            error_msg = f"Error in convert_powerpoint_to_pdf: {str(e)}"
            print(f"Error in convert_powerpoint_to_pdf: {error_msg}")
            print(traceback.format_exc())
            return {"success": False, "error": error_msg}

    def _convert_with_libreoffice(self, pptx_path: str, output_path: str) -> bool:
        """Convert using LibreOffice headless mode (best quality)"""
        try:
            # Check if LibreOffice is available
            result = subprocess.run(['which', 'libreoffice'], 
                                  capture_output=True, text=True)
            if result.returncode != 0:
                print("LibreOffice not found, trying alternative commands")
                # Try alternative LibreOffice commands
                for cmd in ['soffice', 'libreoffice', '/usr/bin/libreoffice']:
                    result = subprocess.run(['which', cmd], 
                                          capture_output=True, text=True)
                    if result.returncode == 0:
                        libreoffice_cmd = cmd
                        break
                else:
                    return False
            else:
                libreoffice_cmd = 'libreoffice'
            
            # Create temporary directory for conversion
            temp_output_dir = os.path.join(self.temp_dir, f"libreoffice_temp_{uuid.uuid4().hex}")
            os.makedirs(temp_output_dir, exist_ok=True)
            
            # Convert with LibreOffice
            cmd = [
                libreoffice_cmd,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_output_dir,
                pptx_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode == 0:
                # Find the generated PDF file
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                temp_pdf_path = os.path.join(temp_output_dir, f"{base_name}.pdf")
                
                if os.path.exists(temp_pdf_path):
                    # Move to final location
                    shutil.move(temp_pdf_path, output_path)
                    # Clean up temp directory
                    try:
                        shutil.rmtree(temp_output_dir)
                    except:
                        pass
                    return True
            
            print(f"LibreOffice conversion failed: {result.stderr}")
            # Clean up temp directory
            try:
                shutil.rmtree(temp_output_dir)
            except:
                pass
            return False
            
        except Exception as e:
            print(f"LibreOffice conversion method failed: {e}")
            return False

    def _convert_with_reportlab_method(self, pptx_path: str, output_path: str) -> bool:
        """Convert using python-pptx parsing and ReportLab PDF generation"""
        try:
            if not (Presentation and SimpleDocTemplate):
                return False
            
            # Load PowerPoint presentation
            prs = Presentation(pptx_path)
            
            # Create PDF document
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            # Custom styles for better formatting
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=12,
                textColor=colors.darkblue
            )
            
            content_style = ParagraphStyle(
                'CustomContent',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=6
            )
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides):
                print(f"Processing slide {slide_num + 1}")
                
                # Add slide title
                story.append(Paragraph(f"Slide {slide_num + 1}", title_style))
                story.append(Spacer(1, 12))
                
                # Process shapes in slide
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text') and shape.text.strip():
                            # Add text content
                            text_content = shape.text.strip()
                            story.append(Paragraph(text_content, content_style))
                            story.append(Spacer(1, 6))
                        
                        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # Handle images
                            self._add_image_to_story(shape, story)
                        
                        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            # Handle tables
                            self._add_table_to_story(shape, story)
                            
                    except Exception as e:
                        print(f"Error processing shape: {e}")
                        continue
                
                # Add page break after each slide (except last)
                if slide_num < len(prs.slides) - 1:
                    story.append(Spacer(1, 24))
            
            # Build PDF
            doc.build(story)
            return True
            
        except Exception as e:
            print(f"ReportLab conversion method failed: {e}")
            return False

    def _add_image_to_story(self, shape, story):
        """Add image from PowerPoint slide to PDF story"""
        try:
            # Extract image data
            image = shape.image
            image_bytes = image.blob
            
            # Save image temporarily
            temp_img_path = os.path.join(self.temp_dir, f"temp_img_{uuid.uuid4().hex}.png")
            with open(temp_img_path, 'wb') as f:
                f.write(image_bytes)
            
            # Add to story with size constraints
            img = RLImage(temp_img_path, width=4*inch, height=3*inch)
            story.append(img)
            story.append(Spacer(1, 12))
            
            # Clean up temp image
            try:
                os.remove(temp_img_path)
            except:
                pass
                
        except Exception as e:
            print(f"Error adding image: {e}")

    def _add_table_to_story(self, shape, story):
        """Add table from PowerPoint slide to PDF story"""
        try:
            table = shape.table
            data = []
            
            # Extract table data
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_data.append(cell_text)
                data.append(row_data)
            
            if data:
                # Create ReportLab table
                pdf_table = Table(data)
                pdf_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(pdf_table)
                story.append(Spacer(1, 12))
                
        except Exception as e:
            print(f"Error adding table: {e}")

    def _convert_with_pypandoc(self, pptx_path: str, output_path: str) -> bool:
        """Convert using pypandoc as fallback"""
        try:
            if not pypandoc:
                return False
            
            # Convert PowerPoint to PDF using pypandoc
            pypandoc.convert_file(
                pptx_path,
                'pdf',
                outputfile=output_path,
                extra_args=['--pdf-engine=pdflatex']
            )
            
            return os.path.exists(output_path)
            
        except Exception as e:
            print(f"Pypandoc conversion method failed: {e}")
            return False

    def _get_pdf_metadata(self, pdf_path: str) -> Dict[str, Any]:
        """Extract metadata from generated PDF"""
        metadata = {
            "file_size": 0,
            "pages": 0,
            "title": "",
            "author": "",
            "subject": "",
            "creator": "PowerPoint to PDF Converter"
        }
        
        try:
            if os.path.exists(pdf_path):
                metadata["file_size"] = os.path.getsize(pdf_path)
                
                # Try to get page count using PyPDF2 if available
                try:
                    import PyPDF2
                    with open(pdf_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        metadata["pages"] = len(pdf_reader.pages)
                        
                        if pdf_reader.metadata:
                            metadata["title"] = pdf_reader.metadata.get('/Title', '')
                            metadata["author"] = pdf_reader.metadata.get('/Author', '')
                            metadata["subject"] = pdf_reader.metadata.get('/Subject', '')
                except:
                    pass
                    
        except Exception as e:
            print(f"Error extracting PDF metadata: {e}")
            
        return metadata

def test_converter():
    """Test function for the PowerPoint to PDF converter"""
    converter = PowerPointToPDFConverter()
    
    # Test with a sample PowerPoint file
    test_pptx = "test_sample.pptx"
    if not os.path.exists(test_pptx):
        print("No test PowerPoint found, skipping test")
        return
    
    result = converter.convert_powerpoint_to_pdf(test_pptx)
    print(f"PowerPoint to PDF converter test result: {result}")

if __name__ == "__main__":
    test_converter()
