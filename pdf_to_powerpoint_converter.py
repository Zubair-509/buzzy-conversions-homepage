
#!/usr/bin/env python3
"""
Enhanced PDF to PowerPoint Converter
Converts PDF files to PowerPoint presentations while preserving layout, content, images, and tables.
Uses multiple methods for optimal conversion quality with advanced formatting preservation.
"""

import os
import tempfile
import traceback
import shutil
from pathlib import Path
from typing import Dict, Any, Optional, Tuple, List
import uuid
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    print("python-pptx successfully imported")
except ImportError as e:
    print(f"Warning: python-pptx not available - {e}")
    Presentation = None

try:
    from pdf2image import convert_from_path
    print("pdf2image successfully imported")
except ImportError as e:
    print(f"Warning: pdf2image not available - {e}")
    convert_from_path = None

try:
    import pdfplumber
    print("pdfplumber successfully imported")
except ImportError as e:
    print(f"Warning: pdfplumber not available - {e}")
    pdfplumber = None

try:
    from PIL import Image, ImageEnhance, ImageFilter
    print("Pillow successfully imported")
except ImportError as e:
    print(f"Warning: Pillow not available - {e}")
    Image = None

try:
    import fitz  # PyMuPDF
    print("PyMuPDF successfully imported")
except ImportError as e:
    print(f"Warning: PyMuPDF not available - {e}")
    fitz = None

class PDFToPowerPointConverter:
    """Enhanced PDF to PowerPoint converter with advanced image and table handling"""
    
    def __init__(self, temp_dir: str = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self.conversion_cache = {}
        
    def convert_pdf_to_powerpoint(self, pdf_path: str, output_path: str = None) -> Dict[str, Any]:
        """
        Convert PDF to PowerPoint with enhanced image and table preservation
        
        Args:
            pdf_path: Path to input PDF file
            output_path: Path for output PowerPoint file (optional)
            
        Returns:
            Dictionary with conversion results
        """
        try:
            if not os.path.exists(pdf_path):
                return {"success": False, "error": f"PDF file not found: {pdf_path}"}
            
            # Generate output path if not provided
            if not output_path:
                base_name = os.path.splitext(os.path.basename(pdf_path))[0]
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.pptx")
            
            print(f"Starting enhanced PDF to PowerPoint conversion: {pdf_path} -> {output_path}")
            
            # Method 1: PyMuPDF-based conversion with images and tables
            if self._convert_with_pymupdf_method(pdf_path, output_path):
                print("PyMuPDF conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "method": "pymupdf_enhanced",
                    "message": "Conversion completed with images, tables, and text preservation"
                }
            
            # Method 2: Advanced conversion with text extraction and images
            if self._convert_with_advanced_method(pdf_path, output_path):
                print("Advanced conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "method": "advanced",
                    "message": "Conversion completed with text and image preservation"
                }
            
            # Method 3: Enhanced image-based conversion (fallback)
            if self._convert_with_enhanced_image_method(pdf_path, output_path):
                print("Enhanced image-based conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "method": "enhanced_image_based",
                    "message": "Conversion completed using enhanced image-based method"
                }
            
            # Method 4: Enhanced text-only conversion (last resort)
            if self._convert_with_enhanced_text_method(pdf_path, output_path):
                print("Enhanced text-based conversion method succeeded")
                return {
                    "success": True, 
                    "output_path": output_path,
                    "method": "enhanced_text_only",
                    "message": "Conversion completed with enhanced text-only method"
                }
            
            return {"success": False, "error": "All conversion methods failed"}
            
        except Exception as e:
            error_msg = f"Error in convert_pdf_to_powerpoint: {str(e)}"
            print(f"Error in convert_pdf_to_powerpoint: {error_msg}")
            print(traceback.format_exc())
            return {"success": False, "error": error_msg}

    def _convert_with_pymupdf_method(self, pdf_path: str, output_path: str) -> bool:
        """Enhanced conversion using PyMuPDF with image and table extraction"""
        try:
            if not (fitz and Presentation):
                return False
            
            # Open PDF with PyMuPDF
            pdf_doc = fitz.open(pdf_path)
            prs = Presentation()
            
            for page_num in range(len(pdf_doc)):
                print(f"Processing page {page_num + 1} with PyMuPDF")
                page = pdf_doc[page_num]
                
                # Add slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Extract and add images
                image_list = page.get_images()
                img_y_offset = 0.5
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            temp_img_path = os.path.join(self.temp_dir, f"page_{page_num}_img_{img_index}.png")
                            
                            with open(temp_img_path, "wb") as img_file:
                                img_file.write(img_data)
                            
                            # Add image to slide
                            try:
                                slide.shapes.add_picture(
                                    temp_img_path,
                                    Inches(0.5 + (img_index % 2) * 5),
                                    Inches(img_y_offset),
                                    Inches(4),
                                    Inches(3)
                                )
                                img_y_offset += 3.5
                            except:
                                pass
                            
                            # Clean up temp image
                            try:
                                os.remove(temp_img_path)
                            except:
                                pass
                        
                        pix = None
                    except Exception as e:
                        print(f"Error processing image {img_index}: {e}")
                        continue
                
                # Extract text and tables
                text_dict = page.get_text("dict")
                self._add_text_and_tables_to_slide(slide, text_dict, page_num)
            
            pdf_doc.close()
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"PyMuPDF conversion method failed: {e}")
            return False

    def _add_text_and_tables_to_slide(self, slide, text_dict: dict, page_num: int):
        """Add extracted text and tables to slide with proper formatting"""
        try:
            blocks = text_dict.get("blocks", [])
            text_y_offset = 0.5
            table_data = []
            
            for block in blocks:
                if "lines" in block:
                    # Process text blocks
                    block_text = ""
                    for line in block["lines"]:
                        line_text = ""
                        for span in line.get("spans", []):
                            line_text += span.get("text", "")
                        if line_text.strip():
                            block_text += line_text + "\n"
                    
                    if block_text.strip():
                        # Check if this looks like a table (multiple columns with aligned text)
                        if self._is_table_like(block_text):
                            table_data.append(block_text)
                        else:
                            # Add as regular text
                            self._add_text_box_to_slide(slide, block_text.strip(), text_y_offset)
                            text_y_offset += 0.8
            
            # Process detected tables
            if table_data:
                self._add_table_to_slide(slide, table_data, text_y_offset)
                
        except Exception as e:
            print(f"Error adding text and tables to slide: {e}")

    def _is_table_like(self, text: str) -> bool:
        """Detect if text block represents a table"""
        lines = text.strip().split('\n')
        if len(lines) < 2:
            return False
        
        # Check for consistent spacing/tabs that indicate columns
        tab_count = text.count('\t')
        space_patterns = [len(re.findall(r'\s{2,}', line)) for line in lines]
        
        return tab_count > 0 or (len(set(space_patterns)) <= 2 and max(space_patterns) > 0)

    def _add_text_box_to_slide(self, slide, text: str, y_offset: float):
        """Add text box to slide with proper formatting"""
        try:
            if len(text) > 500:  # Truncate very long text
                text = text[:500] + "..."
            
            textbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_offset), 
                Inches(9), Inches(0.6)
            )
            text_frame = textbox.text_frame
            text_frame.text = text
            
            # Format text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.name = 'Arial'
            
        except Exception as e:
            print(f"Error adding text box: {e}")

    def _add_table_to_slide(self, slide, table_data: List[str], y_offset: float):
        """Add table to slide from detected table data"""
        try:
            # Parse table data
            all_rows = []
            for table_text in table_data:
                rows = table_text.strip().split('\n')
                for row in rows:
                    if '\t' in row:
                        cols = row.split('\t')
                    else:
                        cols = re.split(r'\s{2,}', row)
                    if len(cols) > 1:
                        all_rows.append([col.strip() for col in cols])
            
            if not all_rows:
                return
            
            # Determine table dimensions
            max_cols = max(len(row) for row in all_rows)
            rows_count = min(len(all_rows), 10)  # Limit table size
            
            # Create table
            table_shape = slide.shapes.add_table(
                rows_count, max_cols,
                Inches(0.5), Inches(y_offset),
                Inches(9), Inches(2)
            )
            table = table_shape.table
            
            # Fill table data
            for i, row_data in enumerate(all_rows[:rows_count]):
                for j, cell_data in enumerate(row_data[:max_cols]):
                    if j < len(row_data):
                        table.cell(i, j).text = str(cell_data)[:50]  # Limit cell text
                        
        except Exception as e:
            print(f"Error adding table to slide: {e}")
    
    def _convert_with_advanced_method(self, pdf_path: str, output_path: str) -> bool:
        """Enhanced advanced conversion combining text extraction with image conversion"""
        try:
            if not (pdfplumber and convert_from_path and Presentation):
                return False
            
            # Create presentation
            prs = Presentation()
            
            # Convert PDF pages to images for background
            try:
                images = convert_from_path(pdf_path, dpi=200, fmt='PNG')
            except Exception as e:
                print(f"Failed to convert PDF to images: {e}")
                return False
            
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, (page, page_image) in enumerate(zip(pdf.pages, images)):
                    print(f"Processing page {page_num + 1} with advanced method")
                    
                    # Add slide
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Enhance and save page image
                    enhanced_image = self._enhance_image(page_image)
                    temp_image_path = os.path.join(self.temp_dir, f"enhanced_page_{page_num}.png")
                    enhanced_image.save(temp_image_path, 'PNG', quality=95, optimize=True)
                    
                    # Add background image
                    slide.shapes.add_picture(
                        temp_image_path, 
                        Inches(0), Inches(0), 
                        Inches(10), Inches(7.5)
                    )
                    
                    # Extract and add tables
                    tables = page.extract_tables()
                    if tables:
                        self._add_extracted_tables_to_slide(slide, tables, page_num)
                    
                    # Extract and add structured text
                    text_elements = page.extract_text_lines()
                    if text_elements:
                        self._add_structured_text_to_slide(slide, text_elements, page_num)
                    
                    # Clean up temp image
                    try:
                        os.remove(temp_image_path)
                    except:
                        pass
            
            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Advanced conversion method failed: {e}")
            return False

    def _enhance_image(self, image):
        """Enhance image quality for better presentation"""
        try:
            # Enhance contrast and sharpness
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)
            
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.1)
            
            return image
        except:
            return image

    def _add_extracted_tables_to_slide(self, slide, tables: List, page_num: int):
        """Add extracted tables to slide with proper formatting"""
        try:
            for table_idx, table in enumerate(tables[:2]):  # Limit to 2 tables per slide
                if not table or len(table) == 0:
                    continue
                
                # Filter out empty rows and columns
                filtered_table = []
                for row in table:
                    if row and any(cell and str(cell).strip() for cell in row):
                        filtered_row = [str(cell).strip() if cell else "" for cell in row]
                        filtered_table.append(filtered_row)
                
                if len(filtered_table) < 1:
                    continue
                
                # Determine table dimensions
                max_cols = max(len(row) for row in filtered_table)
                rows_count = min(len(filtered_table), 8)  # Limit table size
                
                if max_cols > 6:  # Limit columns for readability
                    max_cols = 6
                
                # Create table shape
                y_position = 1.0 + (table_idx * 3.0)
                table_shape = slide.shapes.add_table(
                    rows_count, max_cols,
                    Inches(0.5), Inches(y_position),
                    Inches(9), Inches(2.5)
                )
                ppt_table = table_shape.table
                
                # Fill table data
                for i, row_data in enumerate(filtered_table[:rows_count]):
                    for j in range(max_cols):
                        cell_text = ""
                        if j < len(row_data) and row_data[j]:
                            cell_text = str(row_data[j])[:100]  # Limit cell text
                        
                        cell = ppt_table.cell(i, j)
                        cell.text = cell_text
                        
                        # Format cell
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(9)
                            paragraph.font.name = 'Arial'
                            
        except Exception as e:
            print(f"Error adding extracted tables: {e}")

    def _add_structured_text_to_slide(self, slide, text_elements: List, page_num: int):
        """Add structured text to slide with better positioning"""
        try:
            # Group text elements by proximity and create text boxes
            grouped_text = self._group_text_elements(text_elements)
            
            for group_idx, text_group in enumerate(grouped_text[:3]):  # Limit to 3 text groups
                if not text_group.strip():
                    continue
                
                # Create text box
                y_position = 0.5 + (group_idx * 2.0)
                textbox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(y_position), 
                    Inches(4.5), Inches(1.8)
                )
                text_frame = textbox.text_frame
                text_frame.text = text_group[:300]  # Limit text length
                
                # Format text
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(50, 50, 50)
                    
        except Exception as e:
            print(f"Error adding structured text: {e}")

    def _group_text_elements(self, text_elements: List) -> List[str]:
        """Group text elements by proximity for better organization"""
        grouped = []
        current_group = []
        
        for element in text_elements:
            text = element.get('text', '').strip()
            if text:
                current_group.append(text)
                
                # Start new group if we have enough text or hit certain patterns
                if len(' '.join(current_group)) > 200 or text.endswith(('.', '!', '?', ':')):
                    grouped.append(' '.join(current_group))
                    current_group = []
        
        # Add remaining text
        if current_group:
            grouped.append(' '.join(current_group))
        
        return grouped
    
    def _convert_with_enhanced_image_method(self, pdf_path: str, output_path: str) -> bool:
        """Enhanced image-based conversion with better quality"""
        try:
            if not (convert_from_path and Presentation):
                return False
            
            # Convert PDF pages to high-quality images
            images = convert_from_path(pdf_path, dpi=300, fmt='PNG')
            
            # Create presentation
            prs = Presentation()
            
            for page_num, image in enumerate(images):
                print(f"Converting page {page_num + 1} to enhanced slide")
                
                # Add slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Enhance image quality
                enhanced_image = self._enhance_image(image)
                
                # Resize image to fit slide dimensions optimally
                slide_width, slide_height = 1920, 1080  # 16:9 aspect ratio
                enhanced_image = enhanced_image.resize((slide_width, slide_height), Image.Resampling.LANCZOS)
                
                # Save enhanced image
                temp_image_path = os.path.join(self.temp_dir, f"enhanced_slide_{page_num}.png")
                enhanced_image.save(temp_image_path, 'PNG', quality=95, optimize=True)
                
                # Add image to slide
                slide.shapes.add_picture(
                    temp_image_path, 
                    Inches(0), Inches(0), 
                    Inches(10), Inches(7.5)
                )
                
                # Clean up temp image
                try:
                    os.remove(temp_image_path)
                except:
                    pass
            
            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Enhanced image-based conversion method failed: {e}")
            return False
    
    def _convert_with_enhanced_text_method(self, pdf_path: str, output_path: str) -> bool:
        """Enhanced text-only conversion with better table detection"""
        try:
            if not (pdfplumber and Presentation):
                return False
            
            # Create presentation
            prs = Presentation()
            
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    print(f"Extracting enhanced content from page {page_num + 1}")
                    
                    # Add slide
                    slide_layout = prs.slide_layouts[1]  # Title and Content layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Set title
                    title = slide.shapes.title
                    title.text = f"Page {page_num + 1}"
                    
                    # Try to extract tables first
                    tables = page.extract_tables()
                    if tables:
                        # Use blank layout for tables
                        slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(slide_layout)
                        
                        # Add title manually
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.1), Inches(9), Inches(0.5)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = f"Page {page_num + 1} - Tables and Content"
                        title_para = title_frame.paragraphs[0]
                        title_para.font.size = Pt(18)
                        title_para.font.bold = True
                        
                        # Add tables
                        self._add_extracted_tables_to_slide(slide, tables, page_num)
                    else:
                        # Extract regular text
                        text = page.extract_text()
                        if text and text.strip():
                            # Enhanced text processing
                            processed_text = self._process_text_content(text)
                            
                            # Add content
                            content = slide.placeholders[1]
                            content.text = processed_text[:1500]  # Increased limit
                            
                            # Enhanced text formatting
                            for paragraph in content.text_frame.paragraphs:
                                paragraph.font.size = Pt(12)
                                paragraph.font.name = 'Arial'
                                paragraph.space_after = Pt(6)
                        else:
                            # Add placeholder text
                            content = slide.placeholders[1]
                            content.text = "Content extraction completed, but no readable text was found on this page. The page may contain images, graphics, or non-standard formatting that requires manual review."
            
            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Enhanced text-based conversion method failed: {e}")
            return False

    def _process_text_content(self, text: str) -> str:
        """Process and clean extracted text for better readability"""
        try:
            # Clean up excessive whitespace
            text = re.sub(r'\n\s*\n', '\n\n', text)
            text = re.sub(r'[ \t]+', ' ', text)
            
            # Add bullet points for lists
            lines = text.split('\n')
            processed_lines = []
            
            for line in lines:
                line = line.strip()
                if line:
                    # Detect if line might be a list item
                    if (line.startswith(('•', '-', '*', '1.', '2.', '3.', '4.', '5.')) or
                        re.match(r'^\d+[\.\)]\s', line) or
                        (len(line) < 100 and not line.endswith(('.', '!', '?')))):
                        if not line.startswith('•'):
                            line = '• ' + line
                    processed_lines.append(line)
            
            return '\n'.join(processed_lines)
            
        except Exception as e:
            print(f"Error processing text content: {e}")
            return text
    
    def cleanup_temp_files(self, conversion_id: str):
        """Clean up temporary files for a conversion"""
        try:
            temp_pattern = os.path.join(self.temp_dir, f"*{conversion_id}*")
            import glob
            for file_path in glob.glob(temp_pattern):
                try:
                    os.remove(file_path)
                except:
                    pass
        except Exception as e:
            print(f"Error cleaning up temp files: {e}")

def test_converter():
    """Test function for the enhanced PDF to PowerPoint converter"""
    converter = PDFToPowerPointConverter()
    
    # Test with a sample PDF
    test_pdf = "test_sample.pdf"
    if not os.path.exists(test_pdf):
        print("No test PDF found, skipping test")
        return
    
    result = converter.convert_pdf_to_powerpoint(test_pdf)
    print(f"Enhanced converter test result: {result}")

if __name__ == "__main__":
    test_converter()
